"""Lecture text extractors (Demo B B-F1).

Pure dispatch by file extension with a graceful degradation chain: a missing
binary lib or a corrupt file yields ``ExtractedText(ok=False)`` instead of
raising, so the pipeline can compute an extraction *rate* (B-F1).

Third-party extractors (pypdf/fitz/docx/pptx) are imported lazily inside their
functions: the module imports cleanly with plain Python (unit tests can inject a
fake via ``EXTRACTORS``), and the real libs are only touched on actual files.

Architecture §C4②: ``ExtractorPort`` is the seam; ``extract_path`` / ``extract_dir``
are the pure dispatchers used by ``pipeline.run_demo_b``.
"""
from __future__ import annotations
import os
from typing import Callable, Optional, Protocol, runtime_checkable

from campus.demo_b.types import LectureDoc, ExtractedText, EXTRACTION_RATE_MIN

__all__ = [
    "ExtractorPort", "extract_pdf", "extract_docx", "extract_pptx",
    "extract_md", "extract_txt", "EXTRACTORS", "SUPPORTED_EXTS",
    "lecture_doc", "extract_path", "extract_dir", "extraction_rate",
]


@runtime_checkable
class ExtractorPort(Protocol):
    """path -> ExtractedText. Real libs + injectable fakes both fit."""
    def __call__(self, path: str) -> ExtractedText: ...


def _doc(path: str, ok: bool = True, error: str = "") -> LectureDoc:
    ext = os.path.splitext(path)[1].lstrip(".").lower()
    try:
        size = os.path.getsize(path)
    except OSError:
        size = 0
    return LectureDoc(path=path, ext=ext, size_bytes=size)


def extract_pdf(path: str) -> ExtractedText:
    """PDF: prefer PyMuPDF (fitz), fall back to pypdf. ok=False if both absent/fail."""
    doc = _doc(path)
    try:
        import fitz  # PyMuPDF
        parts: list[str] = []
        with fitz.open(path) as r:
            for page in r:
                parts.append(page.get_text("text") or "")
        return ExtractedText(doc=doc, text="\n".join(parts).strip())
    except Exception:
        pass
    try:
        from pypdf import PdfReader
        parts = []
        for page in PdfReader(path).pages:
            parts.append(page.extract_text() or "")
        return ExtractedText(doc=doc, text="\n".join(parts).strip())
    except Exception as e:
        return ExtractedText(doc=doc, ok=False, error=f"pdf: {e}")


def extract_docx(path: str) -> ExtractedText:
    """DOCX via python-docx (paragraph text)."""
    doc = _doc(path)
    try:
        import docx
        d = docx.Document(path)
        text = "\n".join(p.text for p in d.paragraphs if p.text)
        return ExtractedText(doc=doc, text=text)
    except Exception as e:
        return ExtractedText(doc=doc, ok=False, error=f"docx: {e}")


def extract_pptx(path: str) -> ExtractedText:
    """PPTX via python-pptx (text frames across all slides)."""
    doc = _doc(path)
    try:
        from pptx import Presentation
        prs = Presentation(path)
        parts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    for p in shape.text_frame.paragraphs:
                        t = "".join(r.text for r in p.runs)
                        if t:
                            parts.append(t)
        return ExtractedText(doc=doc, text="\n".join(parts))
    except Exception as e:
        return ExtractedText(doc=doc, ok=False, error=f"pptx: {e}")


def extract_md(path: str) -> ExtractedText:
    """Markdown / plain text via stdlib read."""
    return _read_text(path)


def extract_txt(path: str) -> ExtractedText:
    return _read_text(path)


def _read_text(path: str) -> ExtractedText:
    doc = _doc(path)
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            return ExtractedText(doc=doc, text=f.read())
    except Exception as e:
        return ExtractedText(doc=doc, ok=False, error=f"read: {e}")


# extension -> extractor. Markdown-like extras collapse onto the stdlib reader.
EXTRACTORS: dict[str, Callable[[str], ExtractedText]] = {
    "pdf": extract_pdf,
    "docx": extract_docx,
    "pptx": extract_pptx,
    "ppt": extract_pptx,
    "md": extract_md,
    "markdown": extract_md,
    "txt": extract_txt,
    "text": extract_txt,
}
SUPPORTED_EXTS = tuple(EXTRACTORS.keys())


def lecture_doc(path: str) -> LectureDoc:
    return _doc(path)


def extract_path(path: str, extractors: Optional[dict] = None) -> ExtractedText:
    """Extract one file by extension. Unsupported ext -> ok=False (not a raise).

    ``extractors`` lets tests inject a fake table without touching the real libs.
    """
    table = extractors if extractors is not None else EXTRACTORS
    ext = os.path.splitext(path)[1].lstrip(".").lower()
    fn = table.get(ext)
    if fn is None:
        return ExtractedText(doc=_doc(path), ok=False,
                             error=f"unsupported ext: .{ext}")
    try:
        return fn(path)
    except Exception as e:
        return ExtractedText(doc=_doc(path), ok=False, error=f"{ext}: {e}")


def extract_dir(path: str, extractors: Optional[dict] = None) -> list[ExtractedText]:
    """Recursively scan a directory and extract every supported file (B-F1).

    Unsupported files are skipped (not counted in the denominator). Missing dir
    returns an empty list (pipeline surfaces this as extraction_rate=0).

    Phase 9: if ``path`` is a single file (not a dir), extract just that file —
    previously a single-file path silently yielded ``[]``.
    """
    if os.path.isfile(path):
        r = extract_path(path, extractors)
        return [r] if r.ok else []
    if not os.path.isdir(path):
        return []
    out: list[ExtractedText] = []
    for root, _dirs, files in os.walk(path):
        for name in sorted(files):
            ext = os.path.splitext(name)[1].lstrip(".").lower()
            if ext in (extractors or EXTRACTORS):
                out.append(extract_path(os.path.join(root, name), extractors))
    return out


def extraction_rate(results: list[ExtractedText],
                    minimum: float = EXTRACTION_RATE_MIN) -> tuple[float, bool]:
    """Fraction of extracted texts that succeeded; (rate, rate >= minimum).

    Empty corpus -> rate 0.0 (not 1.0) so an empty scan never accidentally passes.
    """
    if not results:
        return 0.0, False
    ok = sum(1 for r in results if r.ok and r.chars > 0)
    rate = ok / len(results)
    return rate, rate >= minimum
