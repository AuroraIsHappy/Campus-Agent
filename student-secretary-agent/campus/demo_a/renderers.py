"""Pure renderers: proposal MD -> .docx/.pptx/.pdf, budget -> .xlsx (A-F1).

python-docx / python-pptx / openpyxl are lazy-imported inside each function so the
package imports cleanly without them (tests use pytest.importorskip; the
deterministic e2e tolerates their absence via graceful skip in role_turns).
soffice (LibreOffice) is optional for PDF; absent -> skip (return False).
"""
from __future__ import annotations
import os
import shutil


def to_docx(proposal_md: str, path: str) -> str:
    import docx  # python-docx
    d = docx.Document()
    for line in proposal_md.splitlines():
        s = line.rstrip()
        if not s.strip():
            continue
        if s.startswith("# "):
            d.add_heading(s[2:].strip(), level=1)
        elif s.startswith("## "):
            d.add_heading(s[3:].strip(), level=2)
        elif s.startswith("### "):
            d.add_heading(s[4:].strip(), level=3)
        elif s.startswith("- ") or s.startswith("* "):
            d.add_paragraph(s[2:].strip(), style="List Bullet")
        else:
            d.add_paragraph(s)
    d.save(path)
    return path


def to_xlsx(budget_rows: list[dict], path: str) -> str:
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "budget"
    if budget_rows:
        headers = list(budget_rows[0].keys())
        ws.append(headers)
        for row in budget_rows:
            ws.append([row.get(h, "") for h in headers])
    wb.save(path)
    return path


def to_pptx(outline, path: str) -> str:
    """outline = iterable of (title, [bullets])."""
    from pptx import Presentation
    prs = Presentation()
    for item in outline:
        title, bullets = item[0], item[1]
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = str(title)
        body = slide.placeholders[1].text_frame
        for i, b in enumerate(bullets):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = str(b)
    prs.save(path)
    return path


def to_pdf(src_path: str, out_path: str) -> bool:
    """Convert via LibreOffice if available; else skip (return False)."""
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return False
    import subprocess
    out_dir = os.path.dirname(os.path.abspath(out_path)) or "."
    subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf", "--outdir", out_dir, src_path],
        check=False, capture_output=True, timeout=120)
    return os.path.exists(out_path)
